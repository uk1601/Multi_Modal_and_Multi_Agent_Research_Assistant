import json
import pandas as pd
import os
from openai import OpenAI
import requests

import PyPDF2
import docx
import pptx
import wave
import mutagen
import base64
import zipfile

from app.config.settings import settings

client = OpenAI(api_key=settings.OPENAI_API_KEY)


class FileProcessor:
    def __init__(self):
        pass

    def search(self, query: str) -> str:
        """Performs a web search using DuckDuckGo."""
        url = f"https://api.duckduckgo.com/?q={query}&format=json"
        response = requests.get(url)
        results = response.json()
        return results.get('Abstract', 'No results found.')

    def read_image(self, file_path: str) -> str:
        """Reads an image file and returns its base64 encoded string."""
        try:
            with open(file_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except Exception as e:
            return f"Error reading image: {str(e)}"

    def read_excel(self, file_path: str) -> str:
        """Reads an Excel file and returns its content as a JSON string."""
        try:
            df = pd.read_excel(file_path)
            return df.to_json(orient='records')
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def read_csv(self, file_path: str) -> str:
        """Reads a CSV file and returns its content as a JSON string."""
        try:
            df = pd.read_csv(file_path)
            return df.to_json(orient='records')
        except Exception as e:
            return f"Error reading CSV file: {str(e)}"

    def read_zip(self, file_path: str) -> str:
        """Reads a ZIP file, returns both the list of its contents and the content of each file inside."""
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_contents = {}

                # Loop over each file in the ZIP archive
                for file_name in zip_ref.namelist():
                    with zip_ref.open(file_name) as file:
                        try:
                            # Read file contents and decode to utf-8 if text-based
                            file_contents[file_name] = file.read().decode('utf-8')
                        except UnicodeDecodeError:
                            # If the file is binary, return its base64 encoded content
                            file_contents[file_name] = base64.b64encode(file.read()).decode('utf-8')

                # Return as a JSON string: the structure is {filename: content}
                return json.dumps(file_contents, indent=4)
        except Exception as e:
            return f"Error reading ZIP file: {str(e)}"

    def read_pdf(self, file_path: str) -> str:
        """Reads a PDF file and returns its text content."""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def read_json(self, file_path: str) -> str:
        """Reads a JSON file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return json.dumps(json.load(file))
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    def read_python(self, file_path: str) -> str:
        """Reads a Python file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return file.read()
        except Exception as e:
            return f"Error reading Python file: {str(e)}"

    def read_docx(self, file_path: str) -> str:
        """Reads a DOCX file and returns its text content."""
        try:
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    def read_pptx(self, file_path: str) -> str:
        """Reads a PPTX file and returns its text content."""
        try:
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"Error reading PPTX file: {str(e)}"

    def read_audio(self, file_path: str) -> str:
        """Transcribes audio using OpenAI's Whisper API. Returns transcription and metadata if available."""
        try:
            # Transcribe audio
            transcription = ""
            metadata = {}

            with open(file_path, "rb") as audio_file:
                transcription_response = client.audio.transcribe(model="whisper-1", file=audio_file)
                transcription = transcription_response.get("text", "")

            # Extract metadata for MP3 or WAV if needed
            if file_path.endswith('.mp3'):
                try:
                    audio = mutagen.File(file_path)

                except Exception:
                    metadata["error"] = "Failed to extract MP3 metadata."

            elif file_path.endswith('.wav'):
                try:
                    with wave.open(file_path, 'rb') as wav:
                        metadata = {
                            "channels": wav.getnchannels(),
                            "sample_width": wav.getsampwidth(),
                            "framerate": wav.getframerate(),
                            "frames": wav.getnframes()
                        }
                except Exception:
                    metadata["error"] = "Failed to extract WAV metadata."

            metadata["transcription"] = transcription
            return json.dumps(metadata)

        except Exception as e:
            return f"Error processing audio file: {str(e)}"

    def read_pdb(self, file_path: str) -> str:
        """Reads a PDB file and returns its content."""
        try:
            with open(file_path, 'r') as file:
                return file.read()
        except Exception as e:
            return f"Error reading PDB file: {str(e)}"

    def read_txt(self, file_path: str) -> str:
        """Reads a TXT file and returns its content."""
        try:
            if not os.path.exists(file_path):
                return f"Error: File {file_path} does not exist."
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"


# Define the tools
file_processor = FileProcessor()
tools = {
    "Search": file_processor.search,
    "ReadImage": file_processor.read_image,
    "ReadExcel": file_processor.read_excel,
    "ReadCSV": file_processor.read_csv,
    "ReadZIP": file_processor.read_zip,
    "ReadPDF": file_processor.read_pdf,
    "ReadJSON": file_processor.read_json,
    "ReadPython": file_processor.read_python,
    "ReadDOCX": file_processor.read_docx,
    "ReadPPTX": file_processor.read_pptx,
    "ReadAudio": file_processor.read_audio,
    "ReadPDB": file_processor.read_pdb,
    "ReadTXT": file_processor.read_txt
}